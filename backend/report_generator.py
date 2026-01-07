from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import uuid
from datetime import datetime


# ============================================================================
# COLOR PALETTE - Professional Blue Theme
# ============================================================================
COLORS = {
    'primary': (37, 99, 235),      # Blue
    'secondary': (99, 102, 241),   # Indigo
    'accent': (16, 185, 129),      # Green
    'dark': (30, 41, 59),          # Slate dark
    'light': (241, 245, 249),      # Slate light
    'white': (255, 255, 255),
    'warning': (245, 158, 11),     # Amber
}


# ============================================================================
# PDF REPORT GENERATOR
# ============================================================================
class PDFReport(FPDF):
    """Professional PDF Report with styled headers, footers, and sections."""
    
    def __init__(self, title="Data Analysis Report"):
        super().__init__()
        self.title = title
        self.set_auto_page_break(auto=True, margin=20)
    
    def header(self):
        if self.page_no() == 1:
            return  # No header on title page
        self.set_fill_color(*COLORS['primary'])
        self.rect(0, 0, 210, 15, 'F')
        self.set_font('Arial', 'B', 10)
        self.set_text_color(*COLORS['white'])
        self.set_y(5)
        self.cell(0, 5, self.title, 0, 1, 'C')
        self.set_text_color(*COLORS['dark'])
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Page {self.page_no()} | Generated on {datetime.now().strftime("%Y-%m-%d")}', 0, 0, 'C')

    def add_title_page(self):
        """Create a professional title page."""
        self.add_page()
        # Background
        self.set_fill_color(*COLORS['primary'])
        self.rect(0, 0, 210, 297, 'F')
        
        # Title
        self.set_font('Arial', 'B', 36)
        self.set_text_color(*COLORS['white'])
        self.set_y(100)
        self.cell(0, 20, 'Data Analysis Report', 0, 1, 'C')
        
        # Subtitle
        self.set_font('Arial', '', 16)
        self.set_text_color(200, 220, 255)
        self.cell(0, 10, 'Comprehensive Data Insights & Recommendations', 0, 1, 'C')
        
        # Date
        self.set_y(250)
        self.set_font('Arial', 'I', 12)
        self.cell(0, 10, f'Generated: {datetime.now().strftime("%B %d, %Y")}', 0, 1, 'C')
        
    def add_section_header(self, title, number=None):
        """Add a styled section header."""
        self.ln(5)
        self.set_fill_color(*COLORS['primary'])
        header_text = f"{number}. {title}" if number else title
        self.set_font('Arial', 'B', 16)
        self.set_text_color(*COLORS['primary'])
        self.cell(0, 10, header_text, 0, 1, 'L')
        self.set_draw_color(*COLORS['primary'])
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(5)
        self.set_text_color(*COLORS['dark'])
        self.set_font('Arial', '', 11)

    def add_subsection(self, title):
        """Add a subsection header."""
        self.ln(3)
        self.set_font('Arial', 'B', 12)
        self.set_text_color(*COLORS['secondary'])
        self.cell(0, 8, title, 0, 1, 'L')
        self.set_text_color(*COLORS['dark'])
        self.set_font('Arial', '', 11)
        self.ln(2)

    def add_paragraph(self, text):
        """Add a paragraph of text."""
        safe_text = clean_text(text).encode('latin-1', 'replace').decode('latin-1')
        self.multi_cell(0, 6, safe_text)
        self.ln(3)

    def add_bullet_point(self, text):
        """Add a bullet point."""
        safe_text = clean_text(text).encode('latin-1', 'replace').decode('latin-1')
        self.set_x(15)
        self.set_font('Arial', '', 11)
        self.cell(5, 6, chr(149), 0, 0)  # Bullet character
        self.multi_cell(0, 6, safe_text)

    def add_key_value(self, key, value):
        """Add a key-value pair."""
        self.set_font('Arial', 'B', 11)
        self.set_text_color(*COLORS['secondary'])
        self.cell(60, 7, f"{key}:", 0, 0)
        self.set_font('Arial', '', 11)
        self.set_text_color(*COLORS['dark'])
        safe_value = str(value).encode('latin-1', 'replace').decode('latin-1')
        self.cell(0, 7, safe_value, 0, 1)

    def add_image_with_caption(self, img_path, caption="", width=170):
        """Add an image with a caption."""
        if os.path.exists(img_path):
            try:
                # Check if we need a new page
                if self.get_y() > 200:
                    self.add_page()
                self.image(img_path, x=20, w=width)
                if caption:
                    self.ln(3)
                    self.set_font('Arial', 'I', 9)
                    self.set_text_color(100, 100, 100)
                    self.cell(0, 5, caption, 0, 1, 'C')
                    self.set_text_color(*COLORS['dark'])
                self.ln(5)
            except Exception as e:
                self.add_paragraph(f"[Image could not be loaded: {e}]")


def clean_text(text):
    """Remove markdown artifacts from text."""
    if not text:
        return ""
    text = str(text)
    text = text.replace('**', '').replace('*', '')
    text = text.replace('### ', '').replace('## ', '').replace('# ', '')
    text = text.replace('```python', '').replace('```', '')
    text = text.replace('\n\n\n', '\n\n')
    return text.strip()


def create_pdf_report(report_data, image_paths=None, output_filename=None):
    """
    Creates a comprehensive professional PDF report.
    
    Args:
        report_data: Either a string (simple text) or dict with structured sections:
            {
                'executive_summary': str,
                'data_overview': str,
                'key_statistics': str,
                'data_quality': str,
                'distribution_analysis': str,
                'trend_analysis': str,
                'correlation_insights': str,
                'key_findings': list[str],
                'recommendations': list[str],
                'image_descriptions': dict[str, str]  # {filename: description}
            }
        image_paths: List of image file paths
        output_filename: Output filename
        
    Returns:
        str: Path to the generated PDF.
    """
    if not output_filename:
        output_filename = f"report_{uuid.uuid4().hex[:8]}.pdf"
    
    image_paths = image_paths or []
    
    pdf = PDFReport()
    
    # Handle simple string input (backward compatibility)
    if isinstance(report_data, str):
        report_data = {
            'executive_summary': report_data,
            'key_findings': [],
            'recommendations': []
        }
    
    # 1. Title Page
    pdf.add_title_page()
    
    # 2. Executive Summary
    pdf.add_page()
    pdf.add_section_header("Executive Summary", 1)
    pdf.add_paragraph(report_data.get('executive_summary', 'Analysis completed successfully.'))
    
    # 3. Data Overview
    pdf.add_section_header("Data Overview", 2)
    pdf.add_paragraph(report_data.get('data_overview', 'The dataset was analyzed for patterns and insights.'))
    
    # 4. Key Statistics
    pdf.add_section_header("Key Statistics", 3)
    pdf.add_paragraph(report_data.get('key_statistics', 'Statistical analysis was performed on the dataset.'))
    
    # 5. Data Quality Analysis
    pdf.add_page()
    pdf.add_section_header("Data Quality Analysis", 4)
    pdf.add_paragraph(report_data.get('data_quality', 'Data quality was assessed for completeness and accuracy.'))
    
    # 6. Distribution Analysis with images
    pdf.add_section_header("Distribution Analysis", 5)
    pdf.add_paragraph(report_data.get('distribution_analysis', 'The distribution of key variables was analyzed.'))
    
    # Add first half of images
    image_descriptions = report_data.get('image_descriptions', {})
    half = len(image_paths) // 2 + 1
    for img_path in image_paths[:half]:
        filename = os.path.basename(img_path)
        caption = image_descriptions.get(filename, f"Figure: {filename}")
        pdf.add_image_with_caption(img_path, caption)
    
    # 7. Trend Analysis
    pdf.add_page()
    pdf.add_section_header("Trend Analysis", 6)
    pdf.add_paragraph(report_data.get('trend_analysis', 'Trends and patterns were identified in the data.'))
    
    # 8. Correlation Insights with remaining images
    pdf.add_section_header("Correlation Insights", 7)
    pdf.add_paragraph(report_data.get('correlation_insights', 'Relationships between variables were examined.'))
    
    for img_path in image_paths[half:]:
        filename = os.path.basename(img_path)
        caption = image_descriptions.get(filename, f"Figure: {filename}")
        pdf.add_image_with_caption(img_path, caption)
    
    # 9. Key Findings
    pdf.add_page()
    pdf.add_section_header("Key Findings", 8)
    findings = report_data.get('key_findings', [])
    if isinstance(findings, list) and findings:
        for finding in findings:
            pdf.add_bullet_point(finding)
    else:
        pdf.add_paragraph("Key findings from the analysis have been documented above.")
    
    # 10. Recommendations
    pdf.add_section_header("Recommendations & Next Steps", 9)
    recommendations = report_data.get('recommendations', [])
    if isinstance(recommendations, list) and recommendations:
        for rec in recommendations:
            pdf.add_bullet_point(rec)
    else:
        pdf.add_paragraph("Based on the analysis, further investigation of identified patterns is recommended.")
    
    # Create output directory for PDFs
    output_dir = os.path.join(os.getcwd(), "outputs", "pdfs")
    os.makedirs(output_dir, exist_ok=True)
    
    # Save
    output_path = os.path.join(output_dir, output_filename)
    pdf.output(output_path)
    return output_path


# ============================================================================
# POWERPOINT REPORT GENERATOR
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Add a professional title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%B %d, %Y")
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 200, 255)
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content, bullet_points=None):
    """Add a content slide with header and text."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    header.line.fill.background()
    
    # Header text
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Content
    if content:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        cleaned = clean_text(content)
        if len(cleaned) > 800:
            cleaned = cleaned[:800] + "..."
        p.text = cleaned
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
    
    # Bullet points
    if bullet_points:
        y_pos = Inches(1.5) if not content else Inches(4.5)
        for i, point in enumerate(bullet_points[:6]):  # Max 6 points
            bp_box = slide.shapes.add_textbox(Inches(0.7), y_pos + Inches(i * 0.6), Inches(11.5), Inches(0.5))
            tf = bp_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {clean_text(point)}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
    
    return slide


def add_image_slide(prs, title, img_path, description=""):
    """Add a slide with an image and description."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    header.line.fill.background()
    
    # Header text
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Add image
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.2), width=Inches(8))
        except Exception as e:
            print(f"Error adding image: {e}")
    
    # Description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(9.2), Inches(1.5), Inches(3.8), Inches(5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clean_text(description)[:400]
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(*COLORS['dark'])
    
    return slide


def add_closing_slide(prs):
    """Add a thank you / closing slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    bg.line.fill.background()
    
    # Text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(220, 220, 255)
    p.alignment = PP_ALIGN.CENTER


def create_ppt_report(report_data, image_paths=None, output_filename=None, title=None):
    """
    Creates a PowerPoint presentation with DYNAMIC pages based on user content.
    Uses consistent text style and color scheme across all slides.
    
    Args:
        report_data: Either a string (text content) or dict with sections
        image_paths: Optional list of image file paths
        output_filename: Output filename
        title: Optional title for the presentation (from user prompt)
        
    Returns:
        str: Path to the generated PPTX.
    """
    if not output_filename:
        output_filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    
    image_paths = image_paths or []
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Helper function to extract sections from paragraph content
    def parse_content_into_slides(content_text, presentation_title):
        """
        Parse user content intelligently into slides with unique titles.
        Looks for headers, numbered sections, or creates meaningful titles.
        """
        slides_data = []
        content_text = clean_text(str(content_text))
        
        # Try to split by common section markers
        # Look for patterns like: "1.", "Section:", "Topic:", headers with :
        import re
        
        # Pattern 1: Numbered sections (1., 2., etc.)
        numbered_pattern = r'(?:^|\n)(\d+\.?\s*)([^\n:]+[:\n])'
        
        # Pattern 2: Sections with colons (Topic: content)
        colon_pattern = r'(?:^|\n)([A-Z][^:\n]{2,40}):\s*'
        
        # Try to find natural section breaks
        lines = content_text.split('\n')
        current_section_title = None
        current_section_content = []
        slide_index = 1
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if this line looks like a title/header
            is_title = False
            new_title = None
            
            # Numbered section (1. Topic, 2. Topic)
            num_match = re.match(r'^(\d+)[.)\s]+(.+)$', line)
            if num_match:
                is_title = True
                new_title = num_match.group(2).strip()[:45]
            
            # Colon-based section (Topic: content)
            elif ':' in line and len(line.split(':')[0]) < 40:
                parts = line.split(':', 1)
                if len(parts[0]) > 3 and parts[0][0].isupper():
                    is_title = True
                    new_title = parts[0].strip()[:45]
                    # Add the content after colon to current
                    if parts[1].strip():
                        current_section_content.append(parts[1].strip())
            
            # All caps line (likely a header)
            elif line.isupper() and len(line) > 3 and len(line) < 50:
                is_title = True
                new_title = line.title()[:45]
            
            if is_title and new_title:
                # Save previous section if exists
                if current_section_content:
                    slides_data.append({
                        'title': current_section_title or f"Section {slide_index}",
                        'content': current_section_content.copy()
                    })
                    slide_index += 1
                
                current_section_title = new_title
                current_section_content = []
            else:
                # Add as content
                current_section_content.append(line)
        
        # Don't forget the last section
        if current_section_content:
            slides_data.append({
                'title': current_section_title or f"Content",
                'content': current_section_content.copy()
            })
        
        # If no sections were found, split content into chunks with generated titles
        if not slides_data:
            all_lines = [l.strip() for l in lines if l.strip()]
            chunk_size = 8  # 8 points per slide
            for i in range(0, len(all_lines), chunk_size):
                chunk = all_lines[i:i+chunk_size]
                # Generate title from first line or index
                if chunk:
                    first_line = chunk[0][:40].replace('•', '').replace('-', '').strip()
                    title = first_line if first_line and len(first_line) > 5 else f"Key Points {i//chunk_size + 1}"
                    slides_data.append({
                        'title': title,
                        'content': chunk
                    })
        
        return slides_data
    
    # Parse content into sections if string
    if isinstance(report_data, str):
        # Use intelligent parsing
        sections = parse_content_into_slides(report_data, title or "Presentation")
        
        # Generate title from first section if not provided
        if not title and sections:
            title = sections[0].get('title', 'Presentation')[:60]
    else:
        # Convert dict to list of (title, content) tuples
        sections = []
        for key, value in report_data.items():
            if key not in ['image_descriptions', 'title'] and value:
                section_title = key.replace('_', ' ').title()
                if isinstance(value, list):
                    sections.append({'title': section_title, 'content': value})
                else:
                    sections.append({'title': section_title, 'content': str(value)})
        # Use provided title or generate from content
        if not title:
            title = report_data.get('title', 'Presentation')
    
    # Slide 1: Title slide with user-provided or generated title
    add_title_slide(prs, title or "Presentation", "Generated Report")
    
    # Create slides dynamically based on content
    def add_dynamic_slide(prs, slide_title, content):
        """Add a slide with content that fits within boundaries.
        Each point on new line, subpoints with proper indentation."""
        from pptx.enum.text import MSO_AUTO_SIZE
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
        
        # Header bar - consistent primary blue
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
        header.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_title[:50]  # Limit title length
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content area - single text box with multiple paragraphs
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        if isinstance(content, list):
            # Process list items - each as separate paragraph
            first_para = True
            for i, point in enumerate(content[:12]):  # Max 12 points per slide
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                
                point_text = clean_text(str(point))
                
                # Check if this is a subpoint (starts with space, tab, or -)
                is_subpoint = point_text.startswith('  ') or point_text.startswith('\t') or point_text.startswith('- ')
                
                if is_subpoint:
                    # Subpoint - add indentation
                    p.text = f"     ○ {point_text.strip().lstrip('-').strip()}"
                    p.font.size = Pt(14)
                    p.level = 1  # Indent level
                else:
                    # Main point
                    p.text = f"• {point_text[:100]}"
                    p.font.size = Pt(16)
                    p.level = 0
                
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(12)  # Space between points
                
        else:
            # Paragraph content - parse line by line for proper formatting
            content_text = clean_text(str(content))
            lines = content_text.split('\n')
            
            first_para = True
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                    
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                
                # Check for bullet points or subpoints
                is_bullet = line.startswith('•') or line.startswith('-') or line.startswith('*')
                is_subpoint = line.startswith('  ') or line.startswith('\t')
                
                if is_bullet:
                    # Already a bullet - keep it
                    p.text = line if line.startswith('•') else f"• {line.lstrip('-*').strip()}"
                    p.font.size = Pt(16)
                    p.level = 0
                elif is_subpoint:
                    # Subpoint
                    p.text = f"     ○ {line.strip().lstrip('-*').strip()}"
                    p.font.size = Pt(14)
                    p.level = 1
                else:
                    # Regular text
                    p.text = line[:150]
                    p.font.size = Pt(16)
                    p.level = 0
                
                p.font.color.rgb = RGBColor(30, 41, 59)
                p.space_after = Pt(8)  # Space between lines
        
        return slide
    
    # Split content into slides - ensure minimum 10 slides
    slide_count = 1  # Already have title slide
    
    for section in sections:
        # Handle new dict format
        if isinstance(section, dict):
            sec_title = section.get('title', 'Content')
            content = section.get('content', '')
        elif isinstance(section, tuple):
            sec_title, content = section
        else:
            lines = str(section).split('\n')
            sec_title = lines[0][:40] if lines else "Content"
            content = '\n'.join(lines[1:]) if len(lines) > 1 else section
        
        if not content:
            continue
            
        # Split long content into multiple slides
        if isinstance(content, str) and len(content) > 500:
            # Split into chunks by lines first
            lines = content.split('\n')
            chunk_lines = []
            for i in range(0, len(lines), 10):  # 10 lines per slide
                chunk = lines[i:i+10]
                chunk_text = '\n'.join(chunk)
                chunk_title = f"{sec_title}" if i == 0 else f"{sec_title} (cont.)"
                add_dynamic_slide(prs, chunk_title, chunk_text)
                slide_count += 1
        elif isinstance(content, list):
            # Split list into groups of 10 (more points per slide)
            for i in range(0, len(content), 10):
                chunk = content[i:i+10]
                chunk_title = f"{sec_title}" if i == 0 else f"{sec_title} (cont.)"
                add_dynamic_slide(prs, chunk_title, chunk)
                slide_count += 1
        else:
            add_dynamic_slide(prs, sec_title, content)
            slide_count += 1
    
    # Add image slides
    for i, img_path in enumerate(image_paths[:4]):
        if os.path.exists(img_path):
            add_image_slide(prs, f"Visualization {i+1}", img_path)
            slide_count += 1
    
    # Ensure minimum 10 slides by adding summary/filler slides if needed
    default_sections = [
        ("Overview", "This presentation covers the key points provided."),
        ("Summary", "The main topics have been outlined in the previous slides."),
        ("Key Points", "Important information has been highlighted throughout."),
        ("Conclusion", "This concludes the main content of the presentation."),
        ("Next Steps", "Consider reviewing the material for further action items."),
        ("Notes", "Additional notes can be added as needed."),
    ]
    
    filler_idx = 0
    while slide_count < 9 and filler_idx < len(default_sections):  # 9 + closing = 10
        add_dynamic_slide(prs, default_sections[filler_idx][0], default_sections[filler_idx][1])
        slide_count += 1
        filler_idx += 1
    
    # Closing slide
    add_closing_slide(prs)
    
    # Create output directory
    output_dir = os.path.join(os.getcwd(), "outputs", "ppts")
    os.makedirs(output_dir, exist_ok=True)
    
    # Save
    output_path = os.path.join(output_dir, output_filename)
    prs.save(output_path)
    return output_path
